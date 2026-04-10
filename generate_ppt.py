# generate_ppt.py — 졸업작품 발표 PPT 자동 생성
# 실행: python generate_ppt.py

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 색상 팔레트 ─────────────────────────────────────
GREEN       = RGBColor(76, 175, 80)
DARK_GREEN  = RGBColor(46, 125, 50)
WHITE       = RGBColor(255, 255, 255)
BLACK       = RGBColor(26, 26, 26)
GRAY        = RGBColor(100, 116, 139)
LIGHT_GRAY  = RGBColor(241, 245, 249)
RED         = RGBColor(220, 38, 38)
BG_COLOR    = RGBColor(244, 247, 249)

LOGO_PATH = os.path.join(os.path.dirname(__file__), "logo.png")
HAS_LOGO  = os.path.exists(LOGO_PATH)


def set_slide_bg(slide, color=BG_COLOR):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color=WHITE, radius=0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=BLACK, align=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_logo_small(slide, left, top, size=Inches(0.5)):
    if HAS_LOGO:
        slide.shapes.add_picture(LOGO_PATH, left, top, size, size)


# ══════════════════════════════════════════════════════
# 슬라이드 1: 표지
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, DARK_GREEN)

# 로고
if HAS_LOGO:
    slide.shapes.add_picture(LOGO_PATH, Inches(5.5), Inches(1.2), Inches(2.3), Inches(2.3))

# 제목
add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(1),
             "SAFEVIEW", font_size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.8),
             "AI 기반 사각지대 위험 감지 시스템", font_size=28, color=RGBColor(200, 230, 200), align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
             "주차 차량 사각지대 기반 보행자 위험 감지 및 시청각 경고 시스템",
             font_size=16, color=RGBColor(180, 210, 180), align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
             "Capstone Design | 2026학년도 1학기",
             font_size=14, color=RGBColor(160, 200, 160), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# 슬라이드 2: 목차
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo_small(slide, Inches(0.5), Inches(0.3))

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
             "          목차", font_size=32, bold=True, color=BLACK)

items = [
    ("01", "프로젝트 개요", "프로젝트 배경 및 목표"),
    ("02", "프로젝트 필요성", "국내 사각지대 사고 현황 및 유사 시스템 분석"),
    ("03", "시스템 아키텍처", "전체 시스템 구조 및 데이터 흐름"),
    ("04", "기술 스택", "사용 기술 및 선정 이유"),
    ("05", "핵심 기능", "실시간 모니터링, ROI 설정, 이벤트 기록"),
    ("06", "위험 판단 로직", "규칙 기반 위험 상태 판단 알고리즘"),
    ("07", "AI 모델 및 학습 전략", "YOLOv8 활용 및 향후 Fine-tuning 계획"),
    ("08", "시연 결과", "프로토타입 실행 화면 및 탐지 결과"),
    ("09", "향후 계획", "2차 개발 방향 및 개선 사항"),
]

for i, (num, title, desc) in enumerate(items):
    y = Inches(1.4) + Inches(i * 0.62)
    # 번호
    add_text_box(slide, Inches(1.5), y, Inches(0.8), Inches(0.5),
                 num, font_size=20, bold=True, color=GREEN)
    # 제목
    add_text_box(slide, Inches(2.3), y, Inches(4), Inches(0.35),
                 title, font_size=18, bold=True, color=BLACK)
    # 설명
    add_text_box(slide, Inches(2.3), y + Inches(0.3), Inches(8), Inches(0.3),
                 desc, font_size=12, color=GRAY)


# ══════════════════════════════════════════════════════
# 슬라이드 3: 프로젝트 개요
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo_small(slide, Inches(0.5), Inches(0.3))

add_text_box(slide, Inches(1.2), Inches(0.3), Inches(10), Inches(0.7),
             "01  프로젝트 개요", font_size=32, bold=True, color=BLACK)

# 배경
add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(2.5))
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(5), Inches(0.4),
             "📌 프로젝트 배경", font_size=20, bold=True, color=DARK_GREEN)
add_text_box(slide, Inches(1.0), Inches(1.9), Inches(5.2), Inches(1.8),
             "• 생활도로·골목·주차장에서 주차 차량 및 구조물로 인한\n  시야 제한 환경에서 보행자 사고 위험 증가\n\n"
             "• 기존 교통 안전 시스템은 횡단보도·교차로에 집중\n  → 골목·사각지대는 사각지대로 방치",
             font_size=14, color=BLACK)

# 목표
add_shape_bg(slide, Inches(6.8), Inches(1.3), Inches(5.5), Inches(2.5))
add_text_box(slide, Inches(7.0), Inches(1.4), Inches(5), Inches(0.4),
             "🎯 프로젝트 목표", font_size=20, bold=True, color=DARK_GREEN)
add_text_box(slide, Inches(7.0), Inches(1.9), Inches(5.2), Inches(1.8),
             "• CCTV 영상에서 사람과 차량을 실시간 AI 인식\n\n"
             "• 사용자 지정 ROI(관심구역) 내 위험 상황 자동 감지\n\n"
             "• 시각적 경고 + 이벤트 클립 자동 저장\n\n"
             "• 별도 하드웨어 없이 기존 CCTV만으로 구현",
             font_size=14, color=BLACK)

# 기대효과
add_shape_bg(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(2.8))
add_text_box(slide, Inches(1.0), Inches(4.3), Inches(10), Inches(0.4),
             "💡 기대 효과", font_size=20, bold=True, color=DARK_GREEN)

effects = [
    ("저비용 구축", "기존 설치된 CCTV 인프라 재활용\n→ 추가 하드웨어 비용 없음"),
    ("실시간 대응", "위험 감지 즉시 시각 경고\n→ 사고 예방 골든타임 확보"),
    ("증거 보존", "이벤트 전후 15초 클립 자동 저장\n→ 사고 분석 및 증거 활용"),
    ("확장 가능", "소프트웨어 기반 솔루션\n→ 골목, 주차장, 스쿨존 등 다양한 환경 적용"),
]
for i, (title, desc) in enumerate(effects):
    x = Inches(1.0) + Inches(i * 2.9)
    add_text_box(slide, x, Inches(4.9), Inches(2.5), Inches(0.35),
                 f"✅ {title}", font_size=14, bold=True, color=BLACK)
    add_text_box(slide, x, Inches(5.3), Inches(2.6), Inches(1.2),
                 desc, font_size=11, color=GRAY)


# ══════════════════════════════════════════════════════
# 슬라이드 4: 프로젝트 필요성
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo_small(slide, Inches(0.5), Inches(0.3))

add_text_box(slide, Inches(1.2), Inches(0.3), Inches(10), Inches(0.7),
             "02  프로젝트 필요성", font_size=32, bold=True, color=BLACK)

add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(2.2))
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(10), Inches(0.4),
             "🚨 국내 사각지대 보행자 사고 현황", font_size=20, bold=True, color=RED)
add_text_box(slide, Inches(1.0), Inches(1.9), Inches(10), Inches(1.5),
             "• 생활도로 교통사고 사망자 중 보행자 비율 약 60% 이상 (도로교통공단)\n"
             "• 주차 차량에 가려진 어린이 사고는 매년 반복 발생\n"
             "• 골목길·이면도로에는 안전 인프라(신호등, CCTV 경고 시스템) 부재\n"
             "• 기존 교통 안전 시스템은 횡단보도·교차로 위주 → 사각지대는 방치",
             font_size=14, color=BLACK)

add_shape_bg(slide, Inches(0.8), Inches(3.8), Inches(11.5), Inches(3.2))
add_text_box(slide, Inches(1.0), Inches(3.9), Inches(10), Inches(0.4),
             "📊 기존 유사 시스템과의 차별점", font_size=20, bold=True, color=DARK_GREEN)

# 비교 테이블 (텍스트 기반)
add_text_box(slide, Inches(1.0), Inches(4.5), Inches(3.5), Inches(0.3),
             "구분", font_size=13, bold=True, color=GRAY)
add_text_box(slide, Inches(4.5), Inches(4.5), Inches(3.5), Inches(0.3),
             "서울시 횡단보도 시스템", font_size=13, bold=True, color=GRAY)
add_text_box(slide, Inches(8.5), Inches(4.5), Inches(3.5), Inches(0.3),
             "SAFEVIEW (본 프로젝트)", font_size=13, bold=True, color=DARK_GREEN)

rows = [
    ("적용 위치", "횡단보도·교차로", "골목·주차장·이면도로"),
    ("감지 목적", "우회전 보행자 알림", "사각지대 보행자 위험 감지"),
    ("경고 대상", "운전자 (전광판)", "관제 담당자 (모니터링 화면)"),
    ("설치 비용", "전용 하드웨어 + 대형 전광판", "기존 CCTV 재활용 (S/W만)"),
    ("확장성", "설치 장소 한정", "CCTV 있는 곳 어디든 적용"),
]
for i, (cat, existing, ours) in enumerate(rows):
    y = Inches(4.9) + Inches(i * 0.38)
    add_text_box(slide, Inches(1.0), y, Inches(3.5), Inches(0.35), cat, font_size=12, bold=True, color=BLACK)
    add_text_box(slide, Inches(4.5), y, Inches(3.5), Inches(0.35), existing, font_size=12, color=GRAY)
    add_text_box(slide, Inches(8.5), y, Inches(3.5), Inches(0.35), ours, font_size=12, bold=True, color=DARK_GREEN)


# ══════════════════════════════════════════════════════
# 슬라이드 5: 시스템 아키텍처
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo_small(slide, Inches(0.5), Inches(0.3))

add_text_box(slide, Inches(1.2), Inches(0.3), Inches(10), Inches(0.7),
             "03  시스템 아키텍처", font_size=32, bold=True, color=BLACK)

# 파이프라인 흐름
stages = [
    ("📹\nInput", "CCTV 카메라\n영상 파일", Inches(0.5)),
    ("📡\nTransmission", "RTSP 스트림\n로컬 파일 입력", Inches(2.7)),
    ("⚙️\nProcessing", "OpenCV 프레임 획득\nYOLOv8 객체 인식\nROI 판단\n상태 판정", Inches(4.9)),
    ("🖥️\nOutput", "실시간 모니터링 화면\n위험 경고 표시\n이벤트 클립 저장\nCSV 로그 기록", Inches(7.8)),
    ("📋\nStorage", "이벤트 이미지\n영상 클립 (15초)\nCSV 로그", Inches(10.5)),
]

for title, desc, x in stages:
    add_shape_bg(slide, x, Inches(1.5), Inches(2.0), Inches(3.0))
    add_text_box(slide, x + Inches(0.1), Inches(1.6), Inches(1.8), Inches(0.8),
                 title, font_size=14, bold=True, color=DARK_GREEN, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(2.5), Inches(1.8), Inches(1.8),
                 desc, font_size=11, color=BLACK, align=PP_ALIGN.CENTER)

# 화살표 (텍스트로 대체)
for x in [Inches(2.5), Inches(4.7), Inches(7.6), Inches(10.3)]:
    add_text_box(slide, x, Inches(2.5), Inches(0.3), Inches(0.5),
                 "→", font_size=24, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# 폴더 구조
add_shape_bg(slide, Inches(0.5), Inches(4.8), Inches(11.8), Inches(2.3))
add_text_box(slide, Inches(0.7), Inches(4.9), Inches(10), Inches(0.4),
             "📁 프로젝트 구조", font_size=18, bold=True, color=DARK_GREEN)
add_text_box(slide, Inches(0.7), Inches(5.4), Inches(11), Inches(1.5),
             "app.py (진입점)  |  config.py (설정)  |  core/ (핵심 모듈: detector, roi_manager, video_source, danger_logic, event_saver)\n"
             "pages/ (UI: 대시보드, 모니터링, ROI 설정, 이벤트 다시보기)  |  saved_events/ (이벤트 저장)  |  logs/ (CSV 로그)  |  roi_configs/ (ROI 좌표)",
             font_size=12, color=BLACK)


# ══════════════════════════════════════════════════════
# 슬라이드 6: 기술 스택
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo_small(slide, Inches(0.5), Inches(0.3))

add_text_box(slide, Inches(1.2), Inches(0.3), Inches(10), Inches(0.7),
             "04  기술 스택", font_size=32, bold=True, color=BLACK)

techs = [
    ("Python", "메인 개발 언어\n풍부한 AI/영상처리 생태계", "🐍"),
    ("YOLOv8", "실시간 객체 인식 모델\nCOCO 데이터셋 사전학습\nperson/car 탐지", "🤖"),
    ("OpenCV", "영상 프레임 처리\nRTSP 스트림 수신\n바운딩 박스 렌더링", "📹"),
    ("Streamlit", "웹 기반 GUI 프레임워크\n빠른 프로토타이핑\n인터랙티브 위젯", "🖥️"),
    ("RTSP", "실시간 CCTV 스트리밍\nH.264 코덱 지원\n저지연 프레임 전송", "📡"),
    ("Cloudflare\nTunnel", "외부 접속 터널링\n무료 HTTPS 제공\n팀원 원격 접속", "☁️"),
]

for i, (name, desc, icon) in enumerate(techs):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + Inches(col * 4.1)
    y = Inches(1.3) + Inches(row * 2.8)

    add_shape_bg(slide, x, y, Inches(3.7), Inches(2.4))
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.5), Inches(0.5),
                 icon, font_size=24, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.8), y + Inches(0.15), Inches(2.5), Inches(0.4),
                 name, font_size=18, bold=True, color=DARK_GREEN)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.7), Inches(3.2), Inches(1.5),
                 desc, font_size=12, color=BLACK)


# ══════════════════════════════════════════════════════
# 슬라이드 7: 핵심 기능
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo_small(slide, Inches(0.5), Inches(0.3))

add_text_box(slide, Inches(1.2), Inches(0.3), Inches(10), Inches(0.7),
             "05  핵심 기능", font_size=32, bold=True, color=BLACK)

features = [
    ("🎥 실시간 모니터링", "CCTV(RTSP) 또는 영상 파일 입력\nYOLOv8으로 사람/차량 실시간 인식\n바운딩 박스 + ROI 오버레이 표시"),
    ("🗺️ ROI 직접 설정", "영상 프레임 위에 마우스 클릭으로\n관심구역 다각형 직접 그리기\n소스별 ROI 저장/재사용"),
    ("🚨 위험 감지 및 경고", "사람+차 동시 감지 + ROI 내 위치 시\n화면 빨간 테두리 + 경고 문구\n정상→위험 전환 순간 이벤트 발생"),
    ("📋 이벤트 자동 저장", "이벤트 전 5초 + 후 10초 클립 저장\n이벤트 캡처 이미지 저장\nCSV 로그 기록 (시간, 소스, 상태)"),
    ("📅 이벤트 다시보기", "캘린더 기반 날짜별 이벤트 조회\n썸네일 미리보기 + 영상 재생\nH.264 자동 변환 (브라우저 재생)"),
    ("📡 외부 공유 배포", "Cloudflare Tunnel로 외부 접속\n팀원/교수님 링크 공유\n원격 공유 모드 (저대역폭)"),
]

for i, (title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(col * 4.2)
    y = Inches(1.3) + Inches(row * 2.9)

    add_shape_bg(slide, x, y, Inches(3.9), Inches(2.5))
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(3.5), Inches(0.4),
                 title, font_size=16, bold=True, color=DARK_GREEN)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.65), Inches(3.5), Inches(1.7),
                 desc, font_size=12, color=BLACK)


# ══════════════════════════════════════════════════════
# 슬라이드 8: 위험 판단 로직
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo_small(slide, Inches(0.5), Inches(0.3))

add_text_box(slide, Inches(1.2), Inches(0.3), Inches(10), Inches(0.7),
             "06  위험 판단 로직", font_size=32, bold=True, color=BLACK)

# 판단 규칙 테이블
add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(3.5))
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(10), Inches(0.4),
             "규칙 기반 위험 상태 판단", font_size=20, bold=True, color=DARK_GREEN)

# 헤더
for col, (text, x) in enumerate([("감지 조건", Inches(1.2)), ("적용 영역", Inches(5.5)), ("판단 결과", Inches(9.0))]):
    add_text_box(slide, x, Inches(2.0), Inches(3), Inches(0.35), text, font_size=14, bold=True, color=GRAY)

rules = [
    ("사람 + 차량 동시 감지, 사람이 ROI 안에 위치", "ROI 영역 내", "🔴 위험", RED),
    ("사람 + 차량 동시 감지, 사람이 ROI 밖에 위치", "ROI 영역 외", "🟢 정상", GREEN),
    ("사람만 감지", "ROI 내/외", "🟢 정상", GREEN),
    ("차량만 감지", "ROI 내/외", "🟢 정상", GREEN),
    ("ROI 미설정", "-", "⚠️ 판단 불가", RGBColor(180, 130, 0)),
]
for i, (cond, area, result, color) in enumerate(rules):
    y = Inches(2.5) + Inches(i * 0.42)
    add_text_box(slide, Inches(1.2), y, Inches(4.2), Inches(0.38), cond, font_size=12, color=BLACK)
    add_text_box(slide, Inches(5.5), y, Inches(3), Inches(0.38), area, font_size=12, color=GRAY)
    add_text_box(slide, Inches(9.0), y, Inches(2.5), Inches(0.38), result, font_size=13, bold=True, color=color)

# 이벤트 발생 조건
add_shape_bg(slide, Inches(0.8), Inches(5.1), Inches(11.5), Inches(1.8))
add_text_box(slide, Inches(1.0), Inches(5.2), Inches(10), Inches(0.4),
             "⚡ 이벤트 발생 조건", font_size=18, bold=True, color=DARK_GREEN)
add_text_box(slide, Inches(1.0), Inches(5.7), Inches(10), Inches(1.0),
             "• 이전 상태: 정상  →  현재 상태: 위험  (전환 순간에만 이벤트 트리거)\n"
             "• 이벤트 발생 시: 캡처 이미지 즉시 저장 → 이후 10초간 추가 녹화 → 전 5초 + 후 10초 클립 저장\n"
             "• 쿨다운: 15초 (후속 녹화 중복 방지)",
             font_size=13, color=BLACK)


# ══════════════════════════════════════════════════════
# 슬라이드 9: AI 모델 및 학습 전략
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo_small(slide, Inches(0.5), Inches(0.3))

add_text_box(slide, Inches(1.2), Inches(0.3), Inches(10), Inches(0.7),
             "07  AI 모델 및 학습 전략", font_size=32, bold=True, color=BLACK)

# 현재 모델
add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(2.8))
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(5), Inches(0.4),
             "📌 현재 (1차 프로토타입)", font_size=18, bold=True, color=DARK_GREEN)
add_text_box(slide, Inches(1.0), Inches(1.9), Inches(5.2), Inches(2.0),
             "• 모델: YOLOv8n (nano) — Ultralytics\n"
             "• 데이터셋: Microsoft COCO (33만장, 80클래스)\n"
             "• 활용 클래스: person (class 0), car (class 2)\n"
             "• 방식: 사전학습 가중치 그대로 사용\n"
             "• 성능: 실제 CCTV 환경에서 person/car 정상 탐지 확인",
             font_size=13, color=BLACK)

# 향후 계획
add_shape_bg(slide, Inches(6.8), Inches(1.3), Inches(5.5), Inches(2.8))
add_text_box(slide, Inches(7.0), Inches(1.4), Inches(5), Inches(0.4),
             "🔮 향후 (2차 고도화)", font_size=18, bold=True, color=DARK_GREEN)
add_text_box(slide, Inches(7.0), Inches(1.9), Inches(5.2), Inches(2.0),
             "• 자체 데이터 수집: 실제 CCTV 영상 프레임 추출\n"
             "• 라벨링: Roboflow/CVAT 활용 바운딩 박스 태깅\n"
             "• Fine-tuning: YOLOv8n 사전학습 가중치 기반\n  Transfer Learning (50 epoch)\n"
             "• 비교 실험: COCO 모델 vs Fine-tuned 모델\n  mAP 정확도 비교 분석",
             font_size=13, color=BLACK)

# 학습 파이프라인
add_shape_bg(slide, Inches(0.8), Inches(4.4), Inches(11.5), Inches(2.7))
add_text_box(slide, Inches(1.0), Inches(4.5), Inches(10), Inches(0.4),
             "🔄 Fine-tuning 학습 파이프라인", font_size=18, bold=True, color=DARK_GREEN)

pipeline = [
    ("1. 데이터 수집", "자택 CCTV 영상\n프레임 추출\n(100~500장)"),
    ("2. 라벨링", "Roboflow에서\n바운딩 박스\n태깅 작업"),
    ("3. 학습", "YOLOv8 Fine-tuning\nTransfer Learning\n50 epoch"),
    ("4. 평가", "mAP 비교 분석\nCOCO vs Custom\n정확도 검증"),
    ("5. 적용", "Fine-tuned 모델\n프로토타입 교체\n성능 향상 확인"),
]
for i, (title, desc) in enumerate(pipeline):
    x = Inches(1.0) + Inches(i * 2.3)
    add_text_box(slide, x, Inches(5.0), Inches(2.0), Inches(0.35),
                 title, font_size=13, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(5.4), Inches(2.0), Inches(1.3),
                 desc, font_size=11, color=GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# 슬라이드 10: 시연 결과 (스크린샷 자리)
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo_small(slide, Inches(0.5), Inches(0.3))

add_text_box(slide, Inches(1.2), Inches(0.3), Inches(10), Inches(0.7),
             "08  시연 결과", font_size=32, bold=True, color=BLACK)

# 스크린샷 자리 (플레이스홀더)
screens = [
    ("정상 상태 화면", "사람만 / 차만 있을 때\n초록색 바운딩 박스\n🟢 정상 표시"),
    ("위험 감지 화면", "사람+차 동시 + ROI 내\n빨간색 바운딩 박스\n🔴 경고 + 이벤트 저장"),
    ("이벤트 다시보기", "캘린더 날짜 선택\n이벤트 썸네일 목록\n영상 클립 재생"),
]
for i, (title, desc) in enumerate(screens):
    x = Inches(0.5) + Inches(i * 4.3)
    # 이미지 자리
    shape = add_shape_bg(slide, x, Inches(1.3), Inches(3.9), Inches(3.5), color=RGBColor(230, 230, 230))
    add_text_box(slide, x + Inches(0.5), Inches(2.5), Inches(3), Inches(1),
                 "📷 스크린샷\n삽입 위치", font_size=16, color=GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(5.0), Inches(3.5), Inches(0.35),
                 title, font_size=16, bold=True, color=DARK_GREEN, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(5.4), Inches(3.5), Inches(1.2),
                 desc, font_size=12, color=BLACK, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# 슬라이드 11: 향후 계획
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo_small(slide, Inches(0.5), Inches(0.3))

add_text_box(slide, Inches(1.2), Inches(0.3), Inches(10), Inches(0.7),
             "09  향후 계획", font_size=32, bold=True, color=BLACK)

plans = [
    ("2차 개발", [
        "자체 데이터 수집 및 Fine-tuning 학습",
        "야간/역광/우천 환경 데이터 증강 (Data Augmentation)",
        "COCO 모델 vs Fine-tuned 모델 비교 실험",
        "mAP 기준 정량적 성능 평가",
    ]),
    ("기능 고도화", [
        "위험 판단 고도화: 거리·속도 기반 충돌 예측",
        "경고음(사운드 알림) 추가",
        "다중 카메라 동시 모니터링",
        "모바일 앱 푸시 알림 연동",
    ]),
    ("배포 및 확장", [
        "클라우드 서버 배포 (상시 접속)",
        "실제 골목·스쿨존 현장 테스트",
        "지자체·관공서 납품 가능한 패키지화",
        "관제센터 연동 API 개발",
    ]),
]

for i, (category, items) in enumerate(plans):
    x = Inches(0.5) + Inches(i * 4.2)
    add_shape_bg(slide, x, Inches(1.3), Inches(3.9), Inches(5.5))
    add_text_box(slide, x + Inches(0.2), Inches(1.4), Inches(3.5), Inches(0.4),
                 f"📌 {category}", font_size=18, bold=True, color=DARK_GREEN)
    for j, item in enumerate(items):
        add_text_box(slide, x + Inches(0.2), Inches(2.0) + Inches(j * 0.5), Inches(3.4), Inches(0.45),
                     f"• {item}", font_size=12, color=BLACK)


# ══════════════════════════════════════════════════════
# 슬라이드 12: 마무리 (Q&A)
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_GREEN)

if HAS_LOGO:
    slide.shapes.add_picture(LOGO_PATH, Inches(5.5), Inches(1.5), Inches(2.3), Inches(2.3))

add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.8),
             "감사합니다", font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.6),
             "SAFEVIEW — AI 기반 사각지대 위험 감지 시스템", font_size=22, color=RGBColor(200, 230, 200), align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
             "Q & A", font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.4),
             "GitHub: github.com/Mingoon77/SAFE-VIEW", font_size=14, color=RGBColor(180, 220, 180), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# 저장
# ══════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(__file__), "SAFEVIEW_발표자료.pptx")
prs.save(output_path)
print(f"PPT 저장 완료: {output_path}")
print(f"총 {len(prs.slides)}장 슬라이드")
